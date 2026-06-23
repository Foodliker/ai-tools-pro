"""
AI PPT 生成器 Pro - 商业化级AI演示文稿生成平台
运行方式: streamlit run ppt_generator.py
"""

import streamlit as st
import json
import re
import os
import sys
import tempfile
import math
from openai import OpenAI
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 添加父目录到路径以导入 db_utils
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
import db_utils

# ==================== 页面配置 ====================
st.set_page_config(
    page_title="SlideAI Pro - AI PPT生成器",
    page_icon="📊",
)

# ==================== 自定义CSS ====================
st.markdown("""
<style>
    /* ===== 消除顶部留白 ===== */
    header[data-testid="stHeader"] { display: none !important; }
    div[data-testid="stToolbar"] { display: none !important; }
    .block-container { padding-top: 1.5rem !important; }
    body, html { margin: 0 !important; padding: 0 !important; }

    /* ===== 全局背景 - 浅色 ===== */
    .stApp {
        background: #f7f8fc !important;
        min-height: 100vh;
    }

    /* ===== 侧边栏 ===== */
    section[data-testid="stSidebar"] {
        background: #ffffff !important;
        border-right: 1px solid #e8e8ef;
    }
    section[data-testid="stSidebar"] h2 { color: #1a1a2e; }

    /* ===== 主标题 ===== */
    .main-header h1 {
        font-size: 2.8rem; font-weight: 800;
        background: linear-gradient(135deg, #667eea 0%, #764ba2 50%, #f093fb 100%);
        -webkit-background-clip: text; -webkit-text-fill-color: transparent;
        background-clip: text; margin-bottom: 0.3rem;
    }
    .main-header .subtitle {
        color: #999; font-size: 1rem; letter-spacing: 2px;
    }

    /* ===== 按钮 ===== */
    .stButton > button[kind="primary"] {
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%) !important;
        color: white !important; border: none !important; border-radius: 12px !important;
        font-size: 1.1rem !important; font-weight: 700 !important;
        padding: 0.8rem 2rem !important;
        box-shadow: 0 4px 15px rgba(102,126,234,0.3) !important;
    }
    .stButton > button[kind="primary"]:hover {
        transform: translateY(-2px);
        box-shadow: 0 6px 25px rgba(102,126,234,0.5) !important;
    }
    .stDownloadButton > button {
        background: #ffffff !important;
        border: 1px solid #ddd !important;
        color: #333 !important; border-radius: 10px !important;
        font-weight: 600 !important;
    }
    .stDownloadButton > button:hover {
        background: #f0f0f5 !important;
        border-color: #667eea !important;
    }

    /* ===== 结果展示 ===== */
    .result-box {
        background: #ffffff;
        border: 1px solid #e8e8ef;
        border-radius: 16px; padding: 2rem; margin: 1rem 0;
        box-shadow: 0 2px 12px rgba(0,0,0,0.04);
    }
    .stat-card {
        background: #ffffff;
        border: 1px solid #e8e8ef;
        border-radius: 12px; padding: 1rem; text-align: center;
        box-shadow: 0 2px 8px rgba(0,0,0,0.03);
    }
    .stat-card .number {
        font-size: 1.8rem; font-weight: 800;
        background: linear-gradient(135deg, #667eea, #f093fb);
        -webkit-background-clip: text; -webkit-text-fill-color: transparent;
    }
    .stat-card .label { color: #999; font-size: 0.75rem; margin-top: 0.3rem; }

    /* ===== 隐藏默认元素 ===== */
    #MainMenu {visibility: hidden;}
    footer {visibility: hidden;}
    .stDeployButton {visibility: hidden;}
    .stExpander {
        background: #ffffff !important;
        border: 1px solid #e8e8ef !important;
        border-radius: 10px !important;
    }

    /* ===== 移动端适配 ===== */
    @media (max-width: 768px) {
        .block-container { padding: 1rem 0.8rem !important; }
        .main-header h1 { font-size: 2rem !important; }
        .main-header .subtitle { font-size: 0.8rem; letter-spacing: 1px; }
        .stat-card .number { font-size: 1.4rem; }
        .stat-card .label { font-size: 0.65rem; }
        .stat-card { padding: 0.6rem; }
        .stButton > button[kind="primary"] {
            font-size: 0.95rem !important;
            padding: 0.6rem 1rem !important;
        }
        .result-box { padding: 1rem; }
    }
</style>
""", unsafe_allow_html=True)

# ==================== 配色方案 ====================
COLOR_THEMES = {
    "商务蓝": {
        "primary": RGBColor(0x1A, 0x56, 0xDB),
        "secondary": RGBColor(0x3B, 0x82, 0xF6),
        "accent": RGBColor(0x60, 0xA5, 0xFA),
        "bg": RGBColor(0xF0, 0xF4, 0xFF),
        "bg_dark": RGBColor(0x1E, 0x3A, 0x5F),
        "text": RGBColor(0x1E, 0x29, 0x3B),
        "subtext": RGBColor(0x64, 0x74, 0x8B),
        "white": RGBColor(0xFF, 0xFF, 0xFF),
        "light_accent": RGBColor(0xDB, 0xEA, 0xFE),
    },
    "科技紫": {
        "primary": RGBColor(0x7C, 0x3A, 0xED),
        "secondary": RGBColor(0xA7, 0x8B, 0xFA),
        "accent": RGBColor(0xC4, 0xB5, 0xFD),
        "bg": RGBColor(0xF5, 0xF3, 0xFF),
        "bg_dark": RGBColor(0x3B, 0x1F, 0x6E),
        "text": RGBColor(0x1E, 0x1B, 0x4B),
        "subtext": RGBColor(0x6B, 0x72, 0x80),
        "white": RGBColor(0xFF, 0xFF, 0xFF),
        "light_accent": RGBColor(0xED, 0xE9, 0xFE),
    },
    "活力橙": {
        "primary": RGBColor(0xEA, 0x58, 0x0C),
        "secondary": RGBColor(0xFB, 0x92, 0x3C),
        "accent": RGBColor(0xFD, 0xBB, 0x69),
        "bg": RGBColor(0xFF, 0xF7, 0xED),
        "bg_dark": RGBColor(0x6B, 0x2F, 0x04),
        "text": RGBColor(0x43, 0x14, 0x07),
        "subtext": RGBColor(0x78, 0x71, 0x6C),
        "white": RGBColor(0xFF, 0xFF, 0xFF),
        "light_accent": RGBColor(0xFF, 0xED, 0xD5),
    },
    "自然绿": {
        "primary": RGBColor(0x05, 0x96, 0x69),
        "secondary": RGBColor(0x34, 0xD3, 0x99),
        "accent": RGBColor(0x6E, 0xE7, 0xB7),
        "bg": RGBColor(0xEC, 0xFD, 0xF5),
        "bg_dark": RGBColor(0x06, 0x4E, 0x3B),
        "text": RGBColor(0x06, 0x4E, 0x3B),
        "subtext": RGBColor(0x6B, 0x72, 0x80),
        "white": RGBColor(0xFF, 0xFF, 0xFF),
        "light_accent": RGBColor(0xD1, 0xFA, 0xE5),
    },
    "简约黑": {
        "primary": RGBColor(0x1F, 0x29, 0x37),
        "secondary": RGBColor(0x6B, 0x72, 0x80),
        "accent": RGBColor(0x9C, 0xA3, 0xAF),
        "bg": RGBColor(0xF9, 0xFA, 0xFB),
        "bg_dark": RGBColor(0x11, 0x18, 0x27),
        "text": RGBColor(0x11, 0x18, 0x27),
        "subtext": RGBColor(0x9C, 0xA3, 0xAF),
        "white": RGBColor(0xFF, 0xFF, 0xFF),
        "light_accent": RGBColor(0xF3, 0xF4, 0xF6),
    },
    "玫瑰红": {
        "primary": RGBColor(0xBE, 0x18, 0x5D),
        "secondary": RGBColor(0xF4, 0x72, 0xB6),
        "accent": RGBColor(0xF9, 0xA8, 0xD4),
        "bg": RGBColor(0xFD, 0xF2, 0xF8),
        "bg_dark": RGBColor(0x50, 0x07, 0x24),
        "text": RGBColor(0x50, 0x07, 0x24),
        "subtext": RGBColor(0x6B, 0x72, 0x80),
        "white": RGBColor(0xFF, 0xFF, 0xFF),
        "light_accent": RGBColor(0xFC, 0xE7, 0xF3),
    },
    "深海蓝": {
        "primary": RGBColor(0x0E, 0x74, 0x90),
        "secondary": RGBColor(0x22, 0xD3, 0xEE),
        "accent": RGBColor(0x67, 0xE8, 0xF9),
        "bg": RGBColor(0xEC, 0xFE, 0xFF),
        "bg_dark": RGBColor(0x08, 0x33, 0x44),
        "text": RGBColor(0x08, 0x33, 0x44),
        "subtext": RGBColor(0x6B, 0x72, 0x80),
        "white": RGBColor(0xFF, 0xFF, 0xFF),
        "light_accent": RGBColor(0xCF, 0xEC, 0xF7),
    },
    "暗夜金": {
        "primary": RGBColor(0xB8, 0x86, 0x0B),
        "secondary": RGBColor(0xDA, 0xA5, 0x20),
        "accent": RGBColor(0xF0, 0xD0, 0x60),
        "bg": RGBColor(0x1A, 0x1A, 0x2E),
        "bg_dark": RGBColor(0x0F, 0x0F, 0x1A),
        "text": RGBColor(0xFF, 0xD7, 0x00),
        "subtext": RGBColor(0xB0, 0xB0, 0xB0),
        "white": RGBColor(0xFF, 0xFF, 0xFF),
        "light_accent": RGBColor(0x2A, 0x2A, 0x3E),
    },
}

# ==================== 图标映射 ====================
ICON_MAP = {
    "目标": "🎯", "数据": "📊", "创新": "💡", "团队": "👥",
    "增长": "📈", "安全": "🔒", "技术": "⚙️", "全球": "🌍",
    "时间": "⏰", "质量": "✅", "教育": "📚", "合作": "🤝",
    "金融": "💰", "健康": "❤️", "环境": "🌱", "未来": "🚀",
    "策略": "♟️", "市场": "📣", "用户": "👤", "产品": "📦",
    "效率": "⚡", "服务": "🛎️", "品牌": "🏷️", "研究": "🔬",
}

# ==================== PPT 生成函数 ====================

def add_shape(slide, shape_type, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    """添加形状的辅助函数"""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=RGBColor(0, 0, 0), alignment=PP_ALIGN.LEFT,
                font_name="微软雅黑"):
    """添加文本框的辅助函数"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multiline_textbox(slide, left, top, width, height, lines, font_size=14,
                          color=RGBColor(0, 0, 0), line_spacing=1.5, bold_first=False,
                          font_name="微软雅黑"):
    """添加多行文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        if bold_first and i == 0:
            p.font.bold = True
        p.space_after = Pt(font_size * (line_spacing - 1))
    return txBox


def set_slide_bg(slide, color):
    """设置幻灯片背景色"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


# ----- 封面页 -----
def create_title_slide(prs, title, subtitle, theme):
    """创建封面页 - 专业设计，带装饰元素"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, theme["bg_dark"])

    # 左侧装饰条
    add_shape(slide, MSO_SHAPE.RECTANGLE,
              Inches(0), Inches(0), Inches(0.15), Inches(7.5),
              fill_color=theme["primary"])

    # 底部装饰块
    add_shape(slide, MSO_SHAPE.RECTANGLE,
              Inches(0), Inches(6.8), Inches(10), Inches(0.7),
              fill_color=theme["primary"])

    # 右侧装饰圆
    add_shape(slide, MSO_SHAPE.OVAL,
              Inches(7.5), Inches(-1), Inches(4), Inches(4),
              fill_color=None, line_color=theme["secondary"], line_width=2)

    # 右下角小装饰圆
    add_shape(slide, MSO_SHAPE.OVAL,
              Inches(8.5), Inches(5.5), Inches(1.5), Inches(1.5),
              fill_color=None, line_color=theme["accent"], line_width=1)

    # 主标题
    add_textbox(slide, Inches(0.8), Inches(2.0), Inches(7.5), Inches(1.8),
                title, font_size=40, bold=True,
                color=theme["white"], alignment=PP_ALIGN.LEFT)

    # 分隔线
    add_shape(slide, MSO_SHAPE.RECTANGLE,
              Inches(0.8), Inches(3.9), Inches(2), Pt(3),
              fill_color=theme["secondary"])

    # 副标题
    if subtitle:
        add_textbox(slide, Inches(0.8), Inches(4.2), Inches(7.5), Inches(0.8),
                    subtitle, font_size=18,
                    color=theme["accent"], alignment=PP_ALIGN.LEFT)

    # 底部文字
    add_textbox(slide, Inches(0.8), Inches(6.9), Inches(8), Inches(0.5),
                "POWERED BY AI  |  PROFESSIONAL PRESENTATION",
                font_size=10, color=theme["white"], alignment=PP_ALIGN.LEFT)


# ----- 目录页 -----
def create_toc_slide(prs, slides_data, theme):
    """创建目录页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, theme["bg"])

    # 顶部装饰条
    add_shape(slide, MSO_SHAPE.RECTANGLE,
              Inches(0), Inches(0), Inches(10), Inches(0.06),
              fill_color=theme["primary"])

    # 标题
    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.7),
                "CONTENTS", font_size=14, bold=False,
                color=theme["subtext"], alignment=PP_ALIGN.LEFT)
    add_textbox(slide, Inches(0.8), Inches(0.9), Inches(8), Inches(0.8),
                "目录", font_size=32, bold=True,
                color=theme["primary"], alignment=PP_ALIGN.LEFT)

    # 目录项
    cols = min(len(slides_data), 6)
    rows = math.ceil(cols / 2)
    for i, sd in enumerate(slides_data[:6]):
        row = i // 2
        col = i % 2
        x = Inches(0.8 + col * 4.5)
        y = Inches(2.0 + row * 1.7)

        # 序号圆
        add_shape(slide, MSO_SHAPE.OVAL,
                  x, y, Inches(0.6), Inches(0.6),
                  fill_color=theme["primary"])
        add_textbox(slide, x, y + Pt(4), Inches(0.6), Inches(0.5),
                    str(i + 1), font_size=18, bold=True,
                    color=theme["white"], alignment=PP_ALIGN.CENTER)

        # 标题文字
        add_textbox(slide, x + Inches(0.8), y + Pt(2), Inches(3.2), Inches(0.5),
                    sd.get("title", f"第{i+1}部分"), font_size=16, bold=True,
                    color=theme["text"], alignment=PP_ALIGN.LEFT)

        # 描述文字
        subtitle = sd.get("subtitle", "")
        if subtitle:
            add_textbox(slide, x + Inches(0.8), y + Inches(0.4), Inches(3.2), Inches(0.4),
                        subtitle, font_size=11,
                        color=theme["subtext"], alignment=PP_ALIGN.LEFT)


# ----- 内容页：标准列表 -----
def create_standard_slide(prs, slide_data, slide_num, total_slides, theme):
    """标准内容页 - 带序号的列表"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, theme["bg"])

    # 顶部装饰条
    add_shape(slide, MSO_SHAPE.RECTANGLE,
              Inches(0), Inches(0), Inches(10), Inches(0.06),
              fill_color=theme["primary"])

    # 左侧装饰线
    add_shape(slide, MSO_SHAPE.RECTANGLE,
              Inches(0.6), Inches(0.4), Pt(4), Inches(0.5),
              fill_color=theme["primary"])

    # 标题
    title = slide_data.get("title", "")
    add_textbox(slide, Inches(0.8), Inches(0.35), Inches(8), Inches(0.7),
                title, font_size=26, bold=True, color=theme["primary"])

    # 副标题
    subtitle = slide_data.get("subtitle", "")
    y_start = Inches(1.3)
    if subtitle:
        add_textbox(slide, Inches(0.8), Inches(1.1), Inches(8.4), Inches(0.4),
                    subtitle, font_size=13, color=theme["subtext"])
        y_start = Inches(1.6)

    # 内容点 - 卡片式
    points = slide_data.get("points", [])
    for i, point in enumerate(points):
        point_text = point.get("text", str(point)) if isinstance(point, dict) else str(point)
        icon = ""
        if isinstance(point, dict):
            icon = point.get("icon", "")

        y = y_start + Inches(i * 1.1)

        # 卡片背景
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                  Inches(0.8), y, Inches(8.4), Inches(0.9),
                  fill_color=theme["white"], line_color=theme["light_accent"], line_width=1)

        # 左侧序号色块
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                  Inches(0.9), y + Inches(0.15), Inches(0.5), Inches(0.6),
                  fill_color=theme["primary"])
        add_textbox(slide, Inches(0.9), y + Inches(0.2), Inches(0.5), Inches(0.5),
                    str(i + 1), font_size=16, bold=True,
                    color=theme["white"], alignment=PP_ALIGN.CENTER)

        # 文字
        add_textbox(slide, Inches(1.6), y + Inches(0.2), Inches(7.3), Inches(0.5),
                    point_text, font_size=14, color=theme["text"])

    # 页码
    add_textbox(slide, Inches(8.5), Inches(7.0), Inches(1.2), Inches(0.3),
                f"{slide_num} / {total_slides}", font_size=10,
                color=theme["subtext"], alignment=PP_ALIGN.RIGHT)


# ----- 内容页：双栏 -----
def create_two_column_slide(prs, slide_data, slide_num, total_slides, theme):
    """双栏布局内容页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, theme["bg"])

    # 顶部装饰条
    add_shape(slide, MSO_SHAPE.RECTANGLE,
              Inches(0), Inches(0), Inches(10), Inches(0.06),
              fill_color=theme["primary"])

    # 标题
    title = slide_data.get("title", "")
    add_textbox(slide, Inches(0.8), Inches(0.35), Inches(8), Inches(0.7),
                title, font_size=26, bold=True, color=theme["primary"])

    subtitle = slide_data.get("subtitle", "")
    if subtitle:
        add_textbox(slide, Inches(0.8), Inches(1.1), Inches(8.4), Inches(0.4),
                    subtitle, font_size=13, color=theme["subtext"])

    # 分栏
    points = slide_data.get("points", [])
    mid = math.ceil(len(points) / 2)
    left_points = points[:mid]
    right_points = points[mid:]

    for col_idx, col_points in enumerate([left_points, right_points]):
        x_base = Inches(0.8 + col_idx * 4.5)

        for i, point in enumerate(col_points):
            point_text = point.get("text", str(point)) if isinstance(point, dict) else str(point)
            y = Inches(1.8 + i * 1.2)

            # 卡片背景
            add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                      x_base, y, Inches(4.0), Inches(1.0),
                      fill_color=theme["white"], line_color=theme["light_accent"], line_width=1)

            # 顶部装饰线
            add_shape(slide, MSO_SHAPE.RECTANGLE,
                      x_base + Inches(0.15), y + Inches(0.08), Inches(0.8), Pt(3),
                      fill_color=theme["secondary"])

            # 文字
            add_textbox(slide, x_base + Inches(0.15), y + Inches(0.25),
                        Inches(3.7), Inches(0.6),
                        point_text, font_size=13, color=theme["text"])

    # 页码
    add_textbox(slide, Inches(8.5), Inches(7.0), Inches(1.2), Inches(0.3),
                f"{slide_num} / {total_slides}", font_size=10,
                color=theme["subtext"], alignment=PP_ALIGN.RIGHT)


# ----- 内容页：图标+文字 -----
def create_icon_grid_slide(prs, slide_data, slide_num, total_slides, theme):
    """图标网格布局"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, theme["bg"])

    # 顶部装饰条
    add_shape(slide, MSO_SHAPE.RECTANGLE,
              Inches(0), Inches(0), Inches(10), Inches(0.06),
              fill_color=theme["primary"])

    # 标题
    title = slide_data.get("title", "")
    add_textbox(slide, Inches(0.8), Inches(0.35), Inches(8), Inches(0.7),
                title, font_size=26, bold=True, color=theme["primary"])

    subtitle = slide_data.get("subtitle", "")
    if subtitle:
        add_textbox(slide, Inches(0.8), Inches(1.1), Inches(8.4), Inches(0.4),
                    subtitle, font_size=13, color=theme["subtext"])

    # 网格
    points = slide_data.get("points", [])
    cols = min(len(points), 3)
    rows = math.ceil(len(points) / cols)

    for i, point in enumerate(points):
        point_text = point.get("text", str(point)) if isinstance(point, dict) else str(point)
        row = i // cols
        col = i % cols
        x = Inches(0.6 + col * 3.1)
        y = Inches(1.8 + row * 2.6)

        # 卡片
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                  x, y, Inches(2.8), Inches(2.2),
                  fill_color=theme["white"], line_color=theme["light_accent"], line_width=1)

        # 图标圆
        add_shape(slide, MSO_SHAPE.OVAL,
                  x + Inches(0.9), y + Inches(0.2), Inches(1.0), Inches(1.0),
                  fill_color=theme["light_accent"])

        # 找匹配的图标
        icon = "📌"
        for keyword, emoji in ICON_MAP.items():
            if keyword in point_text:
                icon = emoji
                break
        add_textbox(slide, x + Inches(0.9), y + Inches(0.35), Inches(1.0), Inches(0.7),
                    icon, font_size=28, alignment=PP_ALIGN.CENTER)

        # 文字（简短）
        short_text = point_text[:20] + ("..." if len(point_text) > 20 else "")
        add_textbox(slide, x + Inches(0.15), y + Inches(1.35), Inches(2.5), Inches(0.7),
                    short_text, font_size=12, color=theme["text"], alignment=PP_ALIGN.CENTER)

    # 页码
    add_textbox(slide, Inches(8.5), Inches(7.0), Inches(1.2), Inches(0.3),
                f"{slide_num} / {total_slides}", font_size=10,
                color=theme["subtext"], alignment=PP_ALIGN.RIGHT)


# ----- 章节分隔页 -----
def create_section_slide(prs, section_title, section_num, theme):
    """章节分隔页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, theme["primary"])

    # 大装饰圆
    add_shape(slide, MSO_SHAPE.OVAL,
              Inches(-2), Inches(-2), Inches(6), Inches(6),
              fill_color=None, line_color=theme["secondary"], line_width=2)

    # 右下角装饰圆
    add_shape(slide, MSO_SHAPE.OVAL,
              Inches(7), Inches(4), Inches(5), Inches(5),
              fill_color=None, line_color=theme["accent"], line_width=1)

    # 序号
    add_textbox(slide, Inches(0.8), Inches(2.0), Inches(2), Inches(1.2),
                f"0{section_num}", font_size=60, bold=True,
                color=theme["white"], alignment=PP_ALIGN.LEFT)

    # 分隔线
    add_shape(slide, MSO_SHAPE.RECTANGLE,
              Inches(0.8), Inches(3.5), Inches(1.5), Pt(3),
              fill_color=theme["white"])

    # 标题
    add_textbox(slide, Inches(0.8), Inches(3.8), Inches(8), Inches(1.2),
                section_title, font_size=32, bold=True,
                color=theme["white"], alignment=PP_ALIGN.LEFT)


# ----- 感谢页 -----
def create_end_slide(prs, theme):
    """创建感谢页 - 带装饰"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, theme["bg_dark"])

    # 左侧装饰条
    add_shape(slide, MSO_SHAPE.RECTANGLE,
              Inches(0), Inches(0), Inches(0.15), Inches(7.5),
              fill_color=theme["primary"])

    # 大装饰圆
    add_shape(slide, MSO_SHAPE.OVAL,
              Inches(6), Inches(-2), Inches(6), Inches(6),
              fill_color=None, line_color=theme["secondary"], line_width=2)
    add_shape(slide, MSO_SHAPE.OVAL,
              Inches(7.5), Inches(3), Inches(3), Inches(3),
              fill_color=None, line_color=theme["accent"], line_width=1)

    # 主文字
    add_textbox(slide, Inches(0.8), Inches(2.2), Inches(8), Inches(1.5),
                "THANK YOU", font_size=48, bold=True,
                color=theme["white"], alignment=PP_ALIGN.LEFT)

    # 分隔线
    add_shape(slide, MSO_SHAPE.RECTANGLE,
              Inches(0.8), Inches(3.9), Inches(2), Pt(3),
              fill_color=theme["secondary"])

    # 副文字
    add_textbox(slide, Inches(0.8), Inches(4.2), Inches(8), Inches(0.8),
                "感谢观看  |  期待与您的合作",
                font_size=18, color=theme["accent"], alignment=PP_ALIGN.LEFT)

    # 底部
    add_textbox(slide, Inches(0.8), Inches(6.9), Inches(8), Inches(0.5),
                "GENERATED BY AI TOOLS PRO",
                font_size=10, color=theme["white"], alignment=PP_ALIGN.LEFT)


# ==================== AI JSON 解析 ====================
def parse_ai_json(text):
    """解析AI返回的JSON"""
    try:
        return json.loads(text)
    except json.JSONDecodeError:
        pass
    patterns = [r'```json\s*\n(.*?)\n\s*```', r'```\s*\n(.*?)\n\s*```']
    for pattern in patterns:
        match = re.search(pattern, text, re.DOTALL)
        if match:
            try:
                return json.loads(match.group(1).strip())
            except json.JSONDecodeError:
                continue
    start = text.find('{')
    end = text.rfind('}')
    if start != -1 and end != -1 and end > start:
        try:
            return json.loads(text[start:end + 1])
        except json.JSONDecodeError:
            pass
    return None


# ==================== 布局选择逻辑 ====================
def choose_layout(slide_index, total_slides, point_count):
    """根据位置和内容量智能选择布局"""
    if slide_index == 0:
        return "icon_grid" if point_count <= 6 else "standard"
    elif slide_index == total_slides - 1:
        return "standard"
    elif point_count <= 4:
        return "icon_grid"
    elif point_count <= 6:
        return "two_column"
    else:
        return "standard"


# ==================== 主生成函数 ====================
def generate_pptx(ppt_data, theme, output_path):
    """生成PPT文件"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slides = ppt_data.get("slides", [])
    total = len(slides)

    # 1. 封面
    create_title_slide(prs, ppt_data.get("title", ""), ppt_data.get("subtitle", ""), theme)

    # 2. 目录页（如果有3页以上内容）
    if total >= 3:
        create_toc_slide(prs, slides, theme)

    # 3. 内容页
    for i, slide_data in enumerate(slides):
        points = slide_data.get("points", [])
        layout = choose_layout(i, total, len(points))

        if layout == "icon_grid":
            create_icon_grid_slide(prs, slide_data, i + 1, total, theme)
        elif layout == "two_column":
            create_two_column_slide(prs, slide_data, i + 1, total, theme)
        else:
            create_standard_slide(prs, slide_data, i + 1, total, theme)

    # 4. 感谢页
    create_end_slide(prs, theme)

    prs.save(output_path)
    return output_path


# ==================== API 配置（服务端密钥） ====================
api_key = st.secrets.get("api_key", "")
api_base = st.secrets.get("api_base", "https://api.deepseek.com/v1")
model_name = st.secrets.get("model_name", "deepseek-chat")

# ==================== 主界面 ====================

# 返回按钮
nav_col1, nav_col2, _ = st.columns([1, 1, 4])
with nav_col1:
    if st.button("← 返回首页", key="back_home_ppt"):
        st.switch_page("app.py")
with nav_col2:
    if st.button("📋 历史记录", key="view_history_ppt"):
        st.switch_page("pages/history.py")

st.markdown("""
<div class="main-header">
    <h1>SlideAI Pro</h1>
    <div class="subtitle">AI-POWERED PRESENTATION GENERATOR</div>
</div>
""", unsafe_allow_html=True)

stats = db_utils.get_stats()
col1, col2, col3, col4 = st.columns(4)
with col1:
    st.markdown('<div class="stat-card"><div class="number">8</div><div class="label">配色主题</div></div>', unsafe_allow_html=True)
with col2:
    st.markdown(f'<div class="stat-card"><div class="number">{stats["ppt"]}</div><div class="label">已生成</div></div>', unsafe_allow_html=True)
with col3:
    st.markdown('<div class="stat-card"><div class="number">4</div><div class="label">布局类型</div></div>', unsafe_allow_html=True)
with col4:
    st.markdown('<div class="stat-card"><div class="number">.pptx</div><div class="label">标准格式</div></div>', unsafe_allow_html=True)

st.markdown("<br>", unsafe_allow_html=True)

col_left, col_right = st.columns([2, 1])

with col_left:
    topic = st.text_input("PPT 主题", placeholder="例：人工智能发展趋势、2024年度工作总结、新产品发布方案")
    subtitle = st.text_input("副标题（可选）", placeholder="例：面向管理层的战略汇报")
    num_slides = st.slider("幻灯片数量（不含封面封底）", min_value=3, max_value=15, value=6)
    detail = st.text_area("补充说明（可选）", height=100,
                          placeholder="可以补充：目标受众、需要涵盖的要点、希望突出的内容等")

with col_right:
    st.markdown('<div style="color:#555;font-size:0.9rem;margin-bottom:0.5rem;"> 配色方案</div>', unsafe_allow_html=True)
    theme_name = st.selectbox("选择风格", list(COLOR_THEMES.keys()), label_visibility="collapsed")
    theme = COLOR_THEMES[theme_name]

    st.markdown("**预览：**")
    preview_html = '<div style="display:flex;gap:4px;margin-top:4px;">'
    for label, color in [("主色", theme["primary"]), ("辅色", theme["secondary"]),
                         ("背景", theme["bg"]), ("文字", theme["text"])]:
        hex_c = f"{color[0]:02x}{color[1]:02x}{color[2]:02x}"
        preview_html += f'<div style="width:40px;height:40px;background:#{hex_c};border-radius:6px;border:1px solid #ddd;" title="{label}"></div>'
    preview_html += '</div>'
    st.markdown(preview_html, unsafe_allow_html=True)

    st.markdown("---")
    st.markdown(f'<div style="color:#999;font-size:0.85rem;">**模型：** {model_name} | **预计：** {num_slides + 2} 页</div>', unsafe_allow_html=True)

st.markdown("<br>", unsafe_allow_html=True)

# 生成按钮
if st.button("✨ 一键生成 PPT", type="primary", use_container_width=True):
    if not api_key:
        st.error("⚠️ 系统未配置 API Key，请联系管理员")
    elif not topic.strip():
        st.warning("⚠️ 请输入PPT主题")
    else:
        with st.spinner("🤖 AI 正在构思PPT内容..."):
            try:
                client = OpenAI(api_key=api_key, base_url=api_base)

                system_prompt = """你是一位专业的PPT内容策划师。请根据用户提供的主题生成PPT内容。

你必须严格返回以下格式的JSON（不要包含任何其他文字，不要用markdown代码块包裹）：

{
  "title": "演示文稿主标题",
  "subtitle": "副标题或一句话描述",
  "slides": [
    {
      "title": "章节标题（10字以内，精炼有力）",
      "subtitle": "章节副标题（可选，一句话说明，可为空字符串）",
      "points": [
        {"text": "要点内容（每点15-25字，有实质信息，不要空话）", "icon": "可选emoji图标"},
        {"text": "要点内容", "icon": ""}
      ]
    }
  ]
}

要求：
1. 每页3-6个要点，每个要点包含具体信息（不要只写标题词）
2. 内容逻辑清晰，层层递进，从概述到细节
3. 标题简洁有力（10字以内）
4. 要点要有数据、案例或具体描述，避免空洞的概念堆砌
5. icon字段为每个要点匹配一个合适的emoji（如📊🎯💡⚙️🚀等）
6. 只返回JSON，不要任何额外说明文字"""

                user_prompt = f"""请为以下主题生成PPT内容：

主题：{topic}
副标题：{subtitle if subtitle else '无'}
幻灯片数量：{num_slides}页
补充说明：{detail if detail else '无'}

请确保内容丰富、专业、有深度。直接返回JSON。"""

                response = client.chat.completions.create(
                    model=model_name,
                    messages=[
                        {"role": "system", "content": system_prompt},
                        {"role": "user", "content": user_prompt}
                    ],
                    temperature=0.7,
                    max_tokens=4000
                )

                raw_content = response.choices[0].message.content
                ppt_data = parse_ai_json(raw_content)

                if ppt_data is None:
                    st.error("❌ AI 返回的内容无法解析为JSON，请重试")
                    with st.expander("查看AI原始回复（调试用）"):
                        st.code(raw_content)
                else:
                    tmp_dir = tempfile.mkdtemp()
                    safe_topic = re.sub(r'[\\/:*?"<>|]', '_', topic)[:30]
                    output_path = os.path.join(tmp_dir, f"{safe_topic}.pptx")
                    generate_pptx(ppt_data, theme, output_path)

                    st.session_state["ppt_data"] = ppt_data
                    st.session_state["ppt_path"] = output_path

                    # 保存到数据库
                    try:
                        db_utils.save_record(
                            tool_type="ppt",
                            title=f"PPT - {topic[:30]}",
                            content=json.dumps(ppt_data, ensure_ascii=False),
                            extra_data={"topic": topic, "subtitle": subtitle, "theme": theme_name, "num_slides": num_slides}
                        )
                    except Exception:
                        pass

                    st.success("✅ PPT 生成成功！")

            except Exception as e:
                st.error(f"生成失败: {e}")

# 结果展示
if "ppt_data" in st.session_state:
    ppt_data = st.session_state["ppt_data"]
    ppt_path = st.session_state["ppt_path"]

    st.markdown('<div class="result-box">', unsafe_allow_html=True)
    st.subheader("📋 内容预览")
    st.markdown(f"### {ppt_data.get('title', '')}")
    if ppt_data.get("subtitle"):
        st.caption(ppt_data["subtitle"])

    for i, slide in enumerate(ppt_data.get("slides", [])):
        with st.expander(f"第 {i + 1} 页：{slide.get('title', '')}"):
            if slide.get("subtitle"):
                st.markdown(f"*{slide['subtitle']}*")
            for j, point in enumerate(slide.get("points", [])):
                if isinstance(point, dict):
                    icon = point.get("icon", "📌")
                    text = point.get("text", str(point))
                    st.markdown(f"{icon} {text}")
                else:
                    st.markdown(f"📌 {point}")
    st.markdown('</div>', unsafe_allow_html=True)

    st.divider()
    with open(ppt_path, "rb") as f:
        st.download_button(
            "📥 下载 PPT 文件 (.pptx)",
            data=f.read(),
            file_name=f"{re.sub(r'[\\/:*?\"<>|]', '_', ppt_data.get('title', 'presentation'))}.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            type="primary",
            use_container_width=True
        )

# 底部
st.markdown("<br><br>", unsafe_allow_html=True)
st.markdown("""
<div style="text-align:center;color:#bbb;font-size:0.75rem;padding:1rem 0;">
SlideAI Pro v3.0 | Powered by AI | Built with Streamlit + python-pptx
</div>
""", unsafe_allow_html=True)
